from docx import Document as DocxDocument
from pptx import Presentation
from datetime import datetime
import PyPDF2
import openpyxl
from pptx.util import Inches
from openpyxl import Workbook
import xlwt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from app.models.document import (
    Document, 
    GeneralDocument,
    HRDocument,
    ITDocument,
    ReconciliationDocument,
    MarketingDocument,
    TransformationDocument,
    CommunicationDocument,
    InternalOperationDocument,
    LegalDocument,
    AccountDocument,
    PortfolioRiskDocument,
    UnderwriterDocument,
    BusinessOperationDocument,
    ClientExperienceDocument,
    RecoveryDocument,
    ProductDocument,
    SalesDocument,
)
from app.models.database import db
from io import BytesIO
from PIL import Image
from PyPDF2 import PdfReader, PdfWriter
import os
from app.utils.logger import logger
import subprocess
import tempfile


def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
    return text


def extract_text_from_word(docx_path):
    doc = DocxDocument(docx_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text


def extract_text_from_ppt(pptx_path):
    presentation = Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text

def extract_text_from_excel(excel_path):
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            text += row_text + "\n"
    return text


def generate_docx_from_text(text):
    """Generate a valid DOCX file from extracted text."""
    doc = DocxDocument()
    doc.add_paragraph(text)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def generate_pptx_from_text(text):
    """Generate a valid PPTX file from extracted text."""
    prs = Presentation()
    # Use a title and content layout (typically layout 1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Extracted Text"
    slide.placeholders[1].text = text
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def generate_xlsx_from_text(text):
    """Generate a valid XLSX file from extracted text."""
    wb = Workbook()
    ws = wb.active
    ws["A1"] = text
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output

def generate_xls_from_text(text):
    """Generate a valid XLS file from extracted text."""
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Extracted Text")
    ws.write(0, 0, text)
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output

def generate_pdf_from_text(text):
    """Generate a valid PDF file from extracted text."""
    output = BytesIO()
    c = canvas.Canvas(output, pagesize=letter)
    width, height = letter
    # For simplicity, we draw the text at a fixed position.
    c.drawString(100, height - 100, text)
    c.showPage()
    c.save()
    output.seek(0)
    return output


def save_file(user_id, file, file_extension, text):
        
        document = Document(
            user_id=user_id,
            file_name=file.filename,
            file_type=file_extension,
            upload_date=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            content=text
        )
        
        db.session.add(document)
        db.session.commit()
        return document  
    

def compress_doc(file):
    """
    Compresses a file if its size exceeds 1 MB.
    Returns the compressed file or the original file if no compression is needed.
    """
    MAX_FILE_SIZE = 1 * 1024 * 1024  # 1 MB in bytes

    try:
        # Check file size
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)

        logger.info(f"Original file size: {file_size} bytes")

        if file_size <= MAX_FILE_SIZE:
            logger.debug("No compression needed.")
            return file  # No compression needed

        # Get file extension
        file_extension = os.path.splitext(file.filename)[1].lower()
        logger.info(f"File extension: {file_extension}")

        # Compress based on file type
        if file_extension in ['.jpg', '.jpeg', '.png']:
            logger.info("Compressing image file.")
            img = Image.open(file)
            img = img.convert('RGB')  # Convert to RGB for JPEG compatibility
            output = BytesIO()
            img.save(output, format='JPEG', quality=85, optimize=True)  # Adjust quality for compression
            output.seek(0)
            output.filename = file.filename
            logger.info(f"Compressed image size: {output.getbuffer().nbytes} bytes")
            return output

        elif file_extension == '.pdf':
            logger.info("Compressing PDF file using Ghostscript.")
            # Save the original BytesIO content to a temporary file
            file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_input:
                temp_input.write(file.read())
                temp_input_path = temp_input.name

            # Prepare a temporary file for Ghostscript output
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
                temp_output_path = temp_output.name

            # Build the Ghostscript command
            gs_command = [
                "gswin64c",  # Ensure Ghostscript is installed and in PATH
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/screen",  # Use /screen for low quality & small size; adjust as needed
                "-dNOPAUSE",
                "-dQUIET",
                "-dBATCH",
                f"-sOutputFile={temp_output_path}",
                temp_input_path
            ]
            subprocess.run(gs_command, check=True)

            # Read the compressed PDF back into BytesIO
            with open(temp_output_path, 'rb') as f:
                compressed_pdf_data = f.read()
            output = BytesIO(compressed_pdf_data)
            output.seek(0)
            output.filename = file.filename

            # Clean up temporary files
            os.remove(temp_input_path)
            os.remove(temp_output_path)

            logger.info(f"Compressed PDF size: {output.getbuffer().nbytes} bytes")
            return output

        elif file_extension == '.docx':
            logger.info("Compressing Word document.")
            file.seek(0)
            doc = DocxDocument(file)
            for paragraph in doc.paragraphs:
                if len(paragraph.text) > 1000:  # Truncate long paragraphs
                    paragraph.text = paragraph.text[:1000] + "..."
            output = BytesIO()
            doc.save(output)
            output.seek(0)
            output.filename = file.filename
            logger.info(f"Compressed Word document size: {output.getbuffer().nbytes} bytes")
            return output

        elif file_extension in ['.pptx']:
            logger.info("Compressing PowerPoint file.")
            file.seek(0)
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and len(shape.text) > 1000:
                        shape.text = shape.text[:1000] + "..."
            output = BytesIO()
            ppt.save(output)
            output.seek(0)
            output.filename = file.filename
            logger.info(f"Compressed PowerPoint size: {output.getbuffer().nbytes} bytes")
            return output

        elif file_extension in ['.xlsx', '.xls']:
            logger.info("Compressing Excel file.")
            file.seek(0)
            wb = openpyxl.load_workbook(file)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and len(str(cell.value)) > 1000:
                            cell.value = str(cell.value)[:1000] + "..."
            output = BytesIO()
            wb.save(output)
            output.seek(0)
            output.filename = file.filename
            logger.info(f"Compressed Excel size: {output.getbuffer().nbytes} bytes")
            return output

        else:
            logger.info("Unsupported file type. Returning original file.")
            return file

    except Exception as e:
        logger.error(f"Error compressing file: {e}", exc_info=True)
        raise

def save_account_files(user_id, file, file_type, text):
    document = AccountDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_general_files(user_id, file, file_type, text):
    document = GeneralDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_businessOps_files(user_id, file, file_type, text):
    document = BusinessOperationDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_communication_files(user_id, file, file_type, text):
    document = CommunicationDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_hr_files(user_id, file, file_type, text):
    document = HRDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_legal_files(user_id, file, file_type, text):
    document = LegalDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_marketing_files(user_id, file, file_type, text):
    document = MarketingDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_internalOps_files(user_id, file, file_type, text):
    document = InternalOperationDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_product_files(user_id, file, file_type, text):
    document = ProductDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_reconciliation_files(user_id, file, file_type, text):
    document = ReconciliationDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document
def save_recovery_files(user_id, file, file_type, text):
    document = RecoveryDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_it_files(user_id, file, file_type, text):
    document = ITDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_transformation_files(user_id, file, file_type, text):
    document = TransformationDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_client_experience_files(user_id, file, file_type, text):
    document = ClientExperienceDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document


def save_sales_files(user_id, file, file_type, text):
    document = SalesDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_portfolio_files(user_id, file, file_type, text):
    document = PortfolioRiskDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document

def save_underwriter_files(user_id, file, file_type, text):
    document = UnderwriterDocument(
        file_name=file.filename,
        file_type=file_type,
        content=text,
        uploaded_by=user_id
    )
    db.session.add(document)
    db.session.commit()
        
    
    return document
